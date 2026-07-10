import os
from datetime import datetime
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Design System Colors
COLOR_DARK_BG = RGBColor(15, 23, 42)      # Slate-900 (For Title/Closing)
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # Slate-50 (For Content)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # Slate-800
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate-500
COLOR_PRIMARY = RGBColor(79, 70, 229)     # Indigo-600
COLOR_PRIMARY_LIGHT = RGBColor(238, 242, 255) # Indigo-50
COLOR_BORDER = RGBColor(226, 232, 240)    # Slate-200

RAG_COLORS = {
    "Green": RGBColor(16, 185, 129),     # Emerald-500
    "Amber": RGBColor(245, 158, 11),     # Amber-500
    "Yellow": RGBColor(245, 158, 11),    # Amber-500
    "Red": RGBColor(239, 68, 68)         # Rose-500
}

RAG_BG_COLORS = {
    "Green": RGBColor(209, 250, 229),    # Emerald-100
    "Amber": RGBColor(254, 243, 199),    # Amber-100
    "Yellow": RGBColor(254, 243, 199),   # Amber-100
    "Red": RGBColor(254, 226, 226)        # Rose-100
}

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def create_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    """Creates a basic shape with solid fill and borders."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
        
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_slide_header(slide, title_text: str, category_text: str = "PORTFOLIO PERFORMANCE SYNTHESIS"):
    """Adds a standard premium header to content slides."""
    # Category Tracker
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.name = FONT_BODY
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.8), Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_DARK
    
    # Divider Line
    create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.02), fill_color=COLOR_BORDER)

def add_slide_footer(slide, slide_num: int):
    """Adds standard footer to content slides."""
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Zycus Professional Services PMO | Confidential"
    p.font.name = FONT_BODY
    p.font.size = Pt(8.5)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p_num = slide.shapes.add_textbox(Inches(11.8), Inches(7.0), Inches(0.8), Inches(0.3))
    tf_num = p_num.text_frame
    p2 = tf_num.paragraphs[0]
    p2.text = f"Slide {slide_num}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.name = FONT_BODY
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED

def style_text_box(tf):
    """Removes margins and enables word wrap."""
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

class PPTPortfolioGenerator:
    def __init__(self, projects_metrics: List[Dict[str, Any]], output_path: str):
        self.projects = projects_metrics
        self.output_path = output_path
        self.prs = Presentation()
        # Set to 16:9 widescreen
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]

    @staticmethod
    def _project_name(project: Dict[str, Any]) -> str:
        """Return a readable project name for slide copy."""
        return str(project.get("project_name", "Unnamed project"))

    def _portfolio_takeaways(self) -> List[str]:
        """Build executive takeaways from the current portfolio metrics."""
        if not self.projects:
            return ["No project data was available for this reporting cycle."]

        behind_schedule = [p for p in self.projects if p.get("schedule_slippage", 0) > 0]
        highest_risk = max(self.projects, key=lambda p: p.get("schedule_slippage", 0))
        total_blockers = sum(p.get("blockers_count", 0) for p in self.projects)
        total_overdue = sum(p.get("overdue_milestones_count", 0) for p in self.projects)
        total_data_notes = sum(len(p.get("data_quality_notes", [])) for p in self.projects)

        takeaways = [
            (
                f"{len(behind_schedule)} of {len(self.projects)} projects are behind schedule. "
                f"{self._project_name(highest_risk)} has the largest schedule variance "
                f"at {highest_risk.get('schedule_slippage', 0) * 100:.1f}%."
            ),
            (
                f"The portfolio contains {total_blockers} active blockers and "
                f"{total_overdue} overdue milestones requiring delivery follow-up."
            ),
        ]

        if total_data_notes:
            takeaways.append(
                f"Input quality needs attention: {total_data_notes} data-quality notes "
                "were recorded across the analyzed plans."
            )
        else:
            takeaways.append("No data-quality warnings were recorded in the analyzed plans.")

        status_summary = ", ".join(
            f"{status}: {sum(p.get('rag_status') == status for p in self.projects)}"
            for status in ("Red", "Amber", "Green")
        )
        takeaways.append(f"Current portfolio distribution is {status_summary}.")
        return takeaways

    def _risk_themes(self) -> List[Dict[str, str]]:
        """Create two concise, data-backed risk themes for the risk slide."""
        highest_blockers = sorted(
            self.projects,
            key=lambda p: p.get("blockers_count", 0),
            reverse=True,
        )[:2]
        highest_overdue = sorted(
            self.projects,
            key=lambda p: p.get("overdue_milestones_count", 0),
            reverse=True,
        )[:2]

        blocker_text = ", ".join(
            f"{self._project_name(p)} ({p.get('blockers_count', 0)})"
            for p in highest_blockers
            if p.get("blockers_count", 0) > 0
        ) or "No active blockers were reported"
        overdue_text = ", ".join(
            f"{self._project_name(p)} ({p.get('overdue_milestones_count', 0)})"
            for p in highest_overdue
            if p.get("overdue_milestones_count", 0) > 0
        ) or "No overdue milestones were reported"

        comments = sum(p.get("neg_comments_count", 0) for p in self.projects)
        quality_notes = sum(len(p.get("data_quality_notes", [])) for p in self.projects)

        return [
            {
                "title": "THEME 1: DELIVERY EXECUTION PRESSURE",
                "body": (
                    f"Active blockers are concentrated in {blocker_text}.\n"
                    f"Overdue milestones are concentrated in {overdue_text}.\n"
                    "These items should drive the next PMO review and recovery plan."
                ),
            },
            {
                "title": "THEME 2: REPORTING AND STAKEHOLDER SIGNALS",
                "body": (
                    f"The plans contain {comments} negative comment signal(s) and "
                    f"{quality_notes} data-quality warning(s).\n"
                    "Validate owners, dates, statuses, and client dependencies before "
                    "the next reporting cycle."
                ),
            },
        ]

    def _recommendations(self) -> List[Dict[str, str]]:
        """Create prioritized recommendations from portfolio-level risk totals."""
        blocker_project = max(self.projects, key=lambda p: p.get("blockers_count", 0))
        milestone_project = max(
            self.projects,
            key=lambda p: p.get("overdue_milestones_count", 0),
        )
        quality_notes = sum(len(p.get("data_quality_notes", [])) for p in self.projects)

        recommendations = [
            {
                "title": "CONTAIN THE LARGEST BLOCKER CLUSTER",
                "text": (
                    f"Run a focused recovery session for {self._project_name(blocker_project)} "
                    f"({blocker_project.get('blockers_count', 0)} active blockers), with an "
                    "owner and target date for each critical item."
                ),
            },
            {
                "title": "RECOVER OVERDUE MILESTONES",
                "text": (
                    f"Rebaseline the overdue milestone path for {self._project_name(milestone_project)} "
                    f"({milestone_project.get('overdue_milestones_count', 0)} overdue milestones) "
                    "and escalate dependencies that cannot be resolved this week."
                ),
            },
            {
                "title": "IMPROVE INPUT QUALITY",
                "text": (
                    f"Close {quality_notes} data-quality warning(s) before the next cycle "
                    "so leadership decisions are based on complete dates, statuses, and comments."
                    if quality_notes
                    else "Keep the weekly data-quality check in place and preserve the current complete input coverage."
                ),
            },
        ]
        return recommendations

    def generate(self):
        """Constructs all slides in the monthly deck."""
        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        
        # Slide counts
        total_projects = len(self.projects)
        red_cnt = sum(1 for p in self.projects if p["rag_status"] == "Red")
        amber_cnt = sum(1 for p in self.projects if p["rag_status"] == "Amber")
        green_cnt = sum(1 for p in self.projects if p["rag_status"] == "Green")
        
        # 1. Slide 1: Title Slide (Dark Theme)
        self._add_title_slide(total_projects)
        
        # 2. Slide 2: Executive Summary (Light Theme)
        self._add_executive_summary_slide(red_cnt, amber_cnt, green_cnt)
        
        # 3. Slide 3: Portfolio Health Overview (Table & Status Layout)
        self._add_portfolio_overview_slide()
        
        # 4. Slide 4: Schedule and Milestone Trends
        self._add_schedule_trends_slide()
        
        # 5. Slide 5: Emerging Risks and Blockers (Client Readiness & ERP integrations)
        self._add_emerging_risks_slide()
        
        # 6. Slide 6: Recommendations
        self._add_recommendations_slide()
        
        # 7. Slide 7: Next Steps / Closing
        self._add_next_steps_slide()
        
        self.prs.save(self.output_path)
        return self.output_path

    def _add_title_slide(self, count: int):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_DARK_BG)
        
        # Accent Block
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(0.08), Inches(2.8), fill_color=COLOR_PRIMARY)
        
        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.0), Inches(3.0))
        tf = tx_box.text_frame
        style_text_box(tf)
        
        p = tf.paragraphs[0]
        p.text = "PROJECT HEALTH MONTHLY SYNTHESIS"
        p.font.name = FONT_HEADING
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(14)
        
        p2 = tf.add_paragraph()
        p2.text = "Executive Client Delivery & Portfolio Status Dashboard"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_PRIMARY
        p2.space_after = Pt(28)
        
        p3 = tf.add_paragraph()
        p3.text = f"Report Date: {datetime.now().strftime('%B %Y')} Reporting Cycle  |  {count} Projects Analyzed"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_TEXT_MUTED

    def _add_executive_summary_slide(self, red: int, amber: int, green: int):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_LIGHT_BG)
        add_slide_header(slide, "Executive Summary & Portfolio Health")
        add_slide_footer(slide, 2)
        
        # Portfolio Breakdown Stats (Left Card)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.6), Inches(4.5), Inches(5.1), fill_color=COLOR_WHITE, line_color=COLOR_BORDER)
        
        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(4.0), Inches(4.5))
        tf = tx_box.text_frame
        style_text_box(tf)
        
        p = tf.paragraphs[0]
        p.text = "PORTFOLIO HEALTH SCORES"
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(16)
        
        # Metric Rows
        def add_summary_stat(tf, label, count, color):
            p_lbl = tf.add_paragraph()
            p_lbl.text = label
            p_lbl.font.name = FONT_BODY
            p_lbl.font.size = Pt(10)
            p_lbl.font.color.rgb = COLOR_TEXT_MUTED
            p_lbl.space_after = Pt(2)
            
            p_val = tf.add_paragraph()
            p_val.text = f"{count} Project{'s' if count != 1 else ''}"
            p_val.font.name = FONT_HEADING
            p_val.font.size = Pt(22)
            p_val.font.bold = True
            p_val.font.color.rgb = color
            p_val.space_after = Pt(14)
            
        add_summary_stat(tf, "Green Status (Healthy)", green, RAG_COLORS["Green"])
        add_summary_stat(tf, "Amber Status (Moderate Risk)", amber, RAG_COLORS["Amber"])
        add_summary_stat(tf, "Red Status (Critical Intervention Required)", red, RAG_COLORS["Red"])
        
        # Key Takeaways (Right Card)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.6), Inches(7.08), Inches(5.1), fill_color=COLOR_WHITE, line_color=COLOR_BORDER)
        
        tx_box2 = slide.shapes.add_textbox(Inches(5.8), Inches(1.9), Inches(6.5), Inches(4.5))
        tf2 = tx_box2.text_frame
        style_text_box(tf2)
        
        p2 = tf2.paragraphs[0]
        p2.text = "PORTFOLIO LEADERSHIP TAKEAWAYS"
        p2.font.name = FONT_HEADING
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_PRIMARY
        p2.space_after = Pt(14)
        
        takeaways = self._portfolio_takeaways()
        
        for take in takeaways:
            p_t = tf2.add_paragraph()
            p_t.text = f"• {take}"
            p_t.font.name = FONT_BODY
            p_t.font.size = Pt(11)
            p_t.font.color.rgb = COLOR_TEXT_DARK
            p_t.line_spacing = 1.3
            p_t.space_after = Pt(12)

    def _add_portfolio_overview_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_LIGHT_BG)
        add_slide_header(slide, "Portfolio Status Overview")
        add_slide_footer(slide, 3)
        
        # Add table
        rows, cols = len(self.projects) + 1, 6
        left, top, width, height = Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Column widths
        table.columns[0].width = Inches(2.2) # Project Name
        table.columns[1].width = Inches(1.8) # Manager
        table.columns[2].width = Inches(1.2) # RAG Status
        table.columns[3].width = Inches(1.3) # % Complete
        table.columns[4].width = Inches(1.8) # Schedule Health
        table.columns[5].width = Inches(3.53) # Key Risk / Blocker
        
        # Header Row
        headers = ["Project Name", "Project Manager", "RAG Status", "% Complete", "Schedule Health", "Key Risk / Blocker"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_PRIMARY
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_HEADING
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.LEFT
            
        # Data Rows
        for row_idx, proj in enumerate(self.projects):
            r = row_idx + 1
            
            # Name
            cell_name = table.cell(r, 0)
            cell_name.text = proj["project_name"].upper()
            
            # PM
            cell_pm = table.cell(r, 1)
            cell_pm.text = proj["project_manager"]
            
            # RAG Status
            cell_rag = table.cell(r, 2)
            cell_rag.text = proj["rag_status"]
            rag_color = RAG_COLORS.get(proj["rag_status"], COLOR_TEXT_DARK)
            p_rag = cell_rag.text_frame.paragraphs[0]
            p_rag.font.bold = True
            p_rag.font.color.rgb = rag_color
            
            # Completion
            cell_pct = table.cell(r, 3)
            cell_pct.text = f"{proj['pct_complete']*100:.1f}%"
            
            # Schedule Health / Stage
            cell_health = table.cell(r, 4)
            cell_health.text = f"Slippage: {proj['schedule_slippage']*100:.1f}%"
            
            # Key Risks
            cell_risk = table.cell(r, 5)
            if proj["blockers"]:
                cell_risk.text = proj["blockers"][0]["task_name"]
            elif proj["overdue_milestones"]:
                cell_risk.text = f"Overdue: {proj['overdue_milestones'][0]['task_name']}"
            else:
                cell_risk.text = "No critical risk flagged"
                
            # Formatting cells font
            for col_idx in range(cols):
                cell = table.cell(r, col_idx)
                # Background white
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
                p = cell.text_frame.paragraphs[0]
                p.font.name = FONT_BODY
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_TEXT_DARK if col_idx != 2 else rag_color

    def _add_schedule_trends_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_LIGHT_BG)
        add_slide_header(slide, "Schedule and Milestone Trends")
        add_slide_footer(slide, 4)
        
        # Display 2 Column Cards side-by-side representing project completion breakdowns
        col_width = Inches(5.6)
        col_gap = Inches(0.6)
        left_start = Inches(0.75)
        
        for idx, proj in enumerate(self.projects[:2]):
            col_left = left_start + idx * (col_width + col_gap)
            create_shape(slide, MSO_SHAPE.RECTANGLE, col_left, Inches(1.6), col_width, Inches(5.1), fill_color=COLOR_WHITE, line_color=COLOR_BORDER)
            
            # Header color block
            create_shape(slide, MSO_SHAPE.RECTANGLE, col_left, Inches(1.6), col_width, Inches(0.4), fill_color=RAG_COLORS[proj["rag_status"]])
            
            tx_box = slide.shapes.add_textbox(col_left + Inches(0.3), Inches(2.2), col_width - Inches(0.6), Inches(4.2))
            tf = tx_box.text_frame
            style_text_box(tf)
            
            p = tf.paragraphs[0]
            p.text = f"{proj['project_name'].upper()} PROGRESS DETAIL"
            p.font.name = FONT_HEADING
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_DARK
            p.space_after = Pt(12)
            
            metrics_list = [
                f"● **Completion Progress**: {proj['pct_complete']*100:.1f}% physical complete vs {proj['time_elapsed_pct']*100:.1f}% elapsed timeline.",
                f"● **Overdue Milestones**: {proj['overdue_milestones_count']} overdue milestones logged in the active schedule.",
                f"● **Tasks In Progress**: {proj['status_counts'].get('In Progress', 0)} subtasks in active development.",
                f"● **Tasks Not Started**: {proj['status_counts'].get('Not Started', 0)} subtasks queued for subsequent phases.",
                f"● **On Hold Tasks**: {proj['status_counts'].get('On Hold', 0)} tasks flagged as on-hold."
            ]
            
            for m in metrics_list:
                p_m = tf.add_paragraph()
                p_m.text = m
                p_m.font.name = FONT_BODY
                p_m.font.size = Pt(11)
                p_m.font.color.rgb = COLOR_TEXT_DARK
                p_m.line_spacing = 1.3
                p_m.space_after = Pt(10)

    def _add_emerging_risks_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_LIGHT_BG)
        add_slide_header(slide, "Emerging Portfolio Risks & Blocker Analysis")
        add_slide_footer(slide, 5)
        
        themes = self._risk_themes()

        # Left box: delivery execution risks
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.1), fill_color=COLOR_WHITE, line_color=COLOR_BORDER)
        tx_box_l = slide.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.0), Inches(4.5))
        tf_l = tx_box_l.text_frame
        style_text_box(tf_l)
        
        p = tf_l.paragraphs[0]
        p.text = themes[0]["title"]
        p.font.name = FONT_HEADING
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(12)
        
        p_body = tf_l.add_paragraph()
        p_body.text = themes[0]["body"]
        p_body.font.name = FONT_BODY
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = COLOR_TEXT_DARK
        p_body.line_spacing = 1.3
        
        # Right box: reporting and stakeholder risks
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.1), fill_color=COLOR_WHITE, line_color=COLOR_BORDER)
        tx_box_r = slide.shapes.add_textbox(Inches(7.28), Inches(1.9), Inches(5.0), Inches(4.5))
        tf_r = tx_box_r.text_frame
        style_text_box(tf_r)
        
        p_r = tf_r.paragraphs[0]
        p_r.text = themes[1]["title"]
        p_r.font.name = FONT_HEADING
        p_r.font.size = Pt(13)
        p_r.font.bold = True
        p_r.font.color.rgb = COLOR_PRIMARY
        p_r.space_after = Pt(12)
        
        p_body_r = tf_r.add_paragraph()
        p_body_r.text = themes[1]["body"]
        p_body_r.font.name = FONT_BODY
        p_body_r.font.size = Pt(11)
        p_body_r.font.color.rgb = COLOR_TEXT_DARK
        p_body_r.line_spacing = 1.3

    def _add_recommendations_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_LIGHT_BG)
        add_slide_header(slide, "Strategic PMO Recommendations")
        add_slide_footer(slide, 6)
        
        col_width = Inches(3.6)
        col_gap = Inches(0.5)
        left_start = Inches(0.75)
        
        recs = self._recommendations()
        
        for idx, rec in enumerate(recs):
            col_left = left_start + idx * (col_width + col_gap)
            create_shape(slide, MSO_SHAPE.RECTANGLE, col_left, Inches(1.6), col_width, Inches(5.1), fill_color=COLOR_WHITE, line_color=COLOR_BORDER)
            
            # Header block
            create_shape(slide, MSO_SHAPE.RECTANGLE, col_left, Inches(1.6), col_width, Inches(0.8), fill_color=COLOR_PRIMARY_LIGHT)
            
            num_box = slide.shapes.add_textbox(col_left + Inches(0.2), Inches(1.75), col_width - Inches(0.4), Inches(0.5))
            p = num_box.text_frame.paragraphs[0]
            p.text = f"ACTION {idx + 1:02d}"
            p.font.name = FONT_HEADING
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY
            
            content_box = slide.shapes.add_textbox(col_left + Inches(0.25), Inches(2.6), col_width - Inches(0.5), Inches(3.8))
            tf = content_box.text_frame
            style_text_box(tf)
            
            p_title = tf.paragraphs[0]
            p_title.text = rec["title"]
            p_title.font.name = FONT_HEADING
            p_title.font.size = Pt(12)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_TEXT_DARK
            p_title.space_after = Pt(10)
            
            p_text = tf.add_paragraph()
            p_text.text = rec["text"]
            p_text.font.name = FONT_BODY
            p_text.font.size = Pt(11)
            p_text.font.color.rgb = COLOR_TEXT_DARK
            p_text.line_spacing = 1.2

    def _add_next_steps_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill_color=COLOR_DARK_BG)
        
        create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(0.08), Inches(2.8), fill_color=COLOR_PRIMARY)
        
        closing_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.0), Inches(3.0))
        tf = closing_box.text_frame
        style_text_box(tf)
        
        p = tf.paragraphs[0]
        p.text = "NEXT STEPS & ROADMAP TO GREEN"
        p.font.name = FONT_HEADING
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(8)
        
        p2 = tf.add_paragraph()
        p2.text = "1. Initiate Weekly Monitoring Cadence using the Reporting Agent"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_PRIMARY
        p2.space_after = Pt(6)
        
        p3 = tf.add_paragraph()
        p3.text = "2. Conduct Data Completeness audits on incoming sheets (ensure baseline start/finish columns)"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(16)
        p3.font.color.rgb = COLOR_PRIMARY
        p3.space_after = Pt(6)
        
        p4 = tf.add_paragraph()
        top_project = max(self.projects, key=lambda p: p.get("blockers_count", 0))
        p4.text = (
            f"3. Review the recovery plan for {self._project_name(top_project)} "
            f"({top_project.get('blockers_count', 0)} active blockers)"
        )
        p4.font.name = FONT_BODY
        p4.font.size = Pt(16)
        p4.font.color.rgb = COLOR_PRIMARY
        p4.space_after = Pt(20)
        
        p5 = tf.add_paragraph()
        p5.text = "Source: Current weekly project health reports | Prepared for PMO review"
        p5.font.name = FONT_BODY
        p5.font.size = Pt(11)
        p5.font.color.rgb = COLOR_TEXT_MUTED
