"""Generate a human-facing weekly PowerPoint for one project."""

import os
import re
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
INK = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
SURFACE = RGBColor(248, 250, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(226, 232, 240)
INDIGO = RGBColor(79, 70, 229)
INDIGO_LIGHT = RGBColor(238, 242, 255)

RAG_COLORS = {
    "Green": RGBColor(16, 185, 129),
    "Amber": RGBColor(245, 158, 11),
    "Yellow": RGBColor(245, 158, 11),
    "Red": RGBColor(239, 68, 68),
}

RAG_BACKGROUND = {
    "Green": RGBColor(209, 250, 229),
    "Amber": RGBColor(254, 243, 199),
    "Yellow": RGBColor(254, 243, 199),
    "Red": RGBColor(254, 226, 226),
}


def _add_box(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def _add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size: int = 18,
    color=INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _add_bullets(slide, items: List[str], left, top, width, height, *, size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(10)
    return box


def _section(reasoning: str, heading: str, fallback: str = "Not available.") -> str:
    """Extract one validated explanation section from the model/rules output."""
    pattern = rf"###\s*{re.escape(heading)}\s*\n(.*?)(?=\n###\s|\Z)"
    match = re.search(pattern, reasoning or "", flags=re.IGNORECASE | re.DOTALL)
    if not match:
        return fallback
    value = match.group(1)
    value = re.sub(r"\*\*(.*?)\*\*", r"\1", value)
    value = re.sub(r"`([^`]*)`", r"\1", value)
    value = re.sub(r"\s+", " ", value).strip()
    return value or fallback


def _shorten(text: str, limit: int) -> str:
    """Keep presentation copy concise enough to fit its assigned space."""
    return text if len(text) <= limit else f"{text[: limit - 1].rstrip()}…"


def _recommended_actions(status: str) -> List[str]:
    """Return brief, audience-facing actions appropriate to the RAG status."""
    if status == "Red":
        return [
            "Escalate the critical blockers and confirm a named owner, decision, and due date.",
            "Run a recovery review for the overdue milestone and its delivery dependencies.",
            "Re-baseline the plan only after recovery dates and client dependencies are confirmed.",
        ]
    if status in {"Amber", "Yellow"}:
        return [
            "Assign owners and due dates for overdue milestones and open dependencies.",
            "Hold a working session to unblock the highest-priority delivery risks.",
            "Review resource coverage and validate recovery dates before the next weekly update.",
        ]
    return [
        "Maintain the current delivery cadence and review the plan for emerging risks.",
        "Confirm owners and dates for upcoming milestones before the next reporting cycle.",
        "Keep project-plan progress and task-status updates current.",
    ]


class WeeklyPPTGenerator:
    """Create a compact, client-readable weekly project health deck."""

    def __init__(self, metrics: Dict[str, Any], output_path: str):
        self.metrics = metrics
        self.output_path = output_path
        self.presentation = Presentation()
        self.presentation.slide_width = SLIDE_WIDTH
        self.presentation.slide_height = SLIDE_HEIGHT
        self.blank_layout = self.presentation.slide_layouts[6]
        self.status = str(metrics.get("rag_status", "Amber"))
        self.status_color = RAG_COLORS.get(self.status, RAG_COLORS["Amber"])

    def _header(self, slide, title: str, subtitle: str):
        _add_text(slide, "WEEKLY PROJECT HEALTH REPORT", Inches(0.7), Inches(0.35), Inches(5.5), Inches(0.3), size=14, color=INDIGO, bold=True)
        _add_text(slide, title, Inches(0.7), Inches(0.7), Inches(11.8), Inches(0.55), size=32, color=INK, bold=True)
        _add_text(slide, subtitle, Inches(0.7), Inches(1.27), Inches(11.8), Inches(0.35), size=16, color=MUTED)
        _add_box(slide, Inches(0.7), Inches(1.72), Inches(11.93), Inches(0.02), BORDER)

    def _footer(self, slide, number: int):
        _add_text(slide, "Zycus Professional Services PMO | Weekly reporting agent", Inches(0.7), Inches(7.08), Inches(8.5), Inches(0.22), size=12, color=MUTED)
        _add_text(slide, f"{number} / 4", Inches(11.5), Inches(7.08), Inches(1.1), Inches(0.22), size=12, color=MUTED, align=PP_ALIGN.RIGHT)

    def _title_slide(self):
        slide = self.presentation.slides.add_slide(self.blank_layout)
        _add_box(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, NAVY)
        _add_box(slide, Inches(0.82), Inches(1.15), Inches(0.12), Inches(4.35), self.status_color)
        _add_text(slide, "WEEKLY PROJECT HEALTH", Inches(1.25), Inches(1.35), Inches(7.5), Inches(0.45), size=18, color=RGBColor(165, 180, 252), bold=True)
        _add_text(slide, str(self.metrics.get("project_name", "Unnamed project")), Inches(1.25), Inches(1.95), Inches(9.5), Inches(1.1), size=44, color=WHITE, bold=True)
        _add_text(slide, f"Project Manager: {self.metrics.get('project_manager', 'Not provided')}  |  Report date: {self.metrics.get('today_date', 'Not provided')}", Inches(1.25), Inches(3.25), Inches(10.5), Inches(0.45), size=18, color=RGBColor(203, 213, 225))
        badge = _add_box(slide, Inches(1.25), Inches(4.4), Inches(2.45), Inches(0.85), RAG_BACKGROUND.get(self.status, RAG_BACKGROUND["Amber"]), radius=True)
        _add_text(slide, f"{self.status.upper()} STATUS", badge.left, badge.top + Inches(0.13), badge.width, Inches(0.45), size=23, color=self.status_color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, "This deck explains the calculated RAG status and the actions that should follow from it.", Inches(1.25), Inches(5.65), Inches(9.7), Inches(0.45), size=17, color=RGBColor(203, 213, 225))

    def _status_slide(self):
        slide = self.presentation.slides.add_slide(self.blank_layout)
        _add_box(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, SURFACE)
        self._header(slide, "The project is currently in a " + self.status.upper() + " state", "Status is based on the current schedule, milestones, blockers, and project-plan data.")
        _add_box(slide, Inches(0.7), Inches(2.05), Inches(3.2), Inches(4.55), CARD, BORDER, radius=True)
        _add_text(slide, "RAG STATUS", Inches(1.0), Inches(2.35), Inches(2.6), Inches(0.3), size=15, color=MUTED, bold=True)
        _add_text(slide, self.status.upper(), Inches(1.0), Inches(2.78), Inches(2.6), Inches(0.7), size=38, color=self.status_color, bold=True)
        _add_text(slide, f"Calculated score: {self.metrics.get('calculated_score', 0):.1f}%", Inches(1.0), Inches(3.65), Inches(2.6), Inches(0.35), size=18, color=INK, bold=True)
        executive_summary = _shorten(
            _section(self.metrics.get("reasoning", ""), "Executive Summary"),
            135,
        )
        _add_text(slide, executive_summary, Inches(1.0), Inches(4.2), Inches(2.6), Inches(1.55), size=16, color=INK)
        metrics = [
            ("Overall progress", f"{self.metrics.get('pct_complete', 0) * 100:.1f}%"),
            ("Timeline elapsed", f"{self.metrics.get('time_elapsed_pct', 0) * 100:.1f}%"),
            ("Schedule slippage", f"{self.metrics.get('schedule_slippage', 0) * 100:.1f}%"),
            ("Active blockers", str(self.metrics.get("blockers_count", 0))),
            ("Overdue milestones", str(self.metrics.get("overdue_milestones_count", 0))),
            ("Project stage", str(self.metrics.get("project_stage", "Not provided"))),
        ]
        for index, (label, value) in enumerate(metrics):
            col = index % 3
            row = index // 3
            left = Inches(4.25 + col * 2.78)
            top = Inches(2.05 + row * 2.28)
            _add_box(slide, left, top, Inches(2.45), Inches(1.8), CARD, BORDER, radius=True)
            _add_text(slide, label, left + Inches(0.2), top + Inches(0.25), Inches(2.05), Inches(0.4), size=15, color=MUTED, bold=True)
            _add_text(slide, value, left + Inches(0.2), top + Inches(0.82), Inches(2.05), Inches(0.65), size=26, color=INK, bold=True)
        self._footer(slide, 2)

    def _reasoning_slide(self):
        slide = self.presentation.slides.add_slide(self.blank_layout)
        _add_box(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, SURFACE)
        self._header(slide, "Why this status was assigned", "Plain-English reasoning based on schedule, milestones, blockers, and available project-plan signals.")
        _add_box(slide, Inches(0.7), Inches(2.05), Inches(5.75), Inches(4.55), CARD, BORDER, radius=True)
        _add_text(slide, "STATUS EXPLANATION", Inches(1.0), Inches(2.35), Inches(4.9), Inches(0.35), size=15, color=INDIGO, bold=True)
        _add_text(slide, _section(self.metrics.get("reasoning", ""), "Why This Status Was Assigned"), Inches(1.0), Inches(2.85), Inches(5.0), Inches(3.25), size=18, color=INK)
        _add_box(slide, Inches(6.75), Inches(2.05), Inches(5.88), Inches(4.55), INDIGO_LIGHT, radius=True)
        _add_text(slide, "WHAT THIS MEANS", Inches(7.05), Inches(2.35), Inches(4.8), Inches(0.35), size=15, color=INDIGO, bold=True)
        action_summary = _recommended_actions(self.status)[0]
        _add_text(slide, action_summary, Inches(7.05), Inches(2.85), Inches(5.0), Inches(1.2), size=21, color=INK, bold=True)
        _add_text(slide, "Use the next weekly review to confirm ownership, unblock decisions, and recovery dates.", Inches(7.05), Inches(4.75), Inches(5.0), Inches(0.85), size=17, color=INK)
        self._footer(slide, 3)

    def _actions_slide(self):
        slide = self.presentation.slides.add_slide(self.blank_layout)
        _add_box(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, SURFACE)
        self._header(slide, "Risk drivers and recommended actions", "Priorities for the next reporting cycle, with assumptions called out explicitly.")
        risks = []
        for blocker in self.metrics.get("blockers", [])[:3]:
            risks.append(f"Blocker: {blocker.get('task_name', 'Unnamed blocker')}")
        for milestone in self.metrics.get("overdue_milestones", [])[:3]:
            risks.append(f"Overdue milestone: {milestone.get('task_name', 'Unnamed milestone')}")
        if not risks:
            risks.append("No active blocker or overdue milestone was identified in the source plan.")
        actions = _recommended_actions(self.status)
        _add_box(slide, Inches(0.7), Inches(2.05), Inches(5.75), Inches(4.55), CARD, BORDER, radius=True)
        _add_text(slide, "TOP RISK DRIVERS", Inches(1.0), Inches(2.35), Inches(4.9), Inches(0.35), size=15, color=self.status_color, bold=True)
        _add_bullets(slide, risks, Inches(1.0), Inches(2.9), Inches(5.0), Inches(3.0), size=18)
        _add_box(slide, Inches(6.75), Inches(2.05), Inches(5.88), Inches(4.55), CARD, BORDER, radius=True)
        _add_text(slide, "RECOMMENDED ACTIONS", Inches(7.05), Inches(2.35), Inches(5.0), Inches(0.35), size=15, color=INDIGO, bold=True)
        _add_bullets(slide, actions, Inches(7.05), Inches(2.9), Inches(5.0), Inches(2.55), size=17)
        assumptions = _shorten(
            _section(self.metrics.get("reasoning", ""), "Data Quality & Assumptions"),
            95,
        )
        _add_text(slide, "Data quality and assumptions", Inches(7.05), Inches(5.55), Inches(4.8), Inches(0.3), size=15, color=MUTED, bold=True)
        _add_text(slide, assumptions, Inches(7.05), Inches(5.9), Inches(5.0), Inches(0.5), size=14, color=INK)
        self._footer(slide, 4)

    def generate(self) -> str:
        os.makedirs(os.path.dirname(self.output_path) or ".", exist_ok=True)
        self._title_slide()
        self._status_slide()
        self._reasoning_slide()
        self._actions_slide()
        self.presentation.save(self.output_path)
        return self.output_path
