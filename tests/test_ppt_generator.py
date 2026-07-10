import os
import tempfile
import unittest

from pptx import Presentation

from src.ppt_generator import PPTPortfolioGenerator


def project_metrics(name, blockers, overdue, slippage, rag="Amber"):
    return {
        "project_name": name,
        "project_manager": "Test Manager",
        "rag_status": rag,
        "pct_complete": 0.5,
        "time_elapsed_pct": 0.6,
        "schedule_slippage": slippage,
        "blockers_count": blockers,
        "overdue_milestones_count": overdue,
        "overdue_milestones": [],
        "blockers": [],
        "status_counts": {"In Progress": 1, "Not Started": 1, "On Hold": 0},
        "neg_comments_count": 2,
        "data_quality_notes": ["Example warning"] if name == "Alpha Delivery" else [],
    }


class TestPPTPortfolioGenerator(unittest.TestCase):
    def setUp(self):
        self.projects = [
            project_metrics("Alpha Delivery", 8, 3, 0.18, "Red"),
            project_metrics("Beta Modernization", 1, 0, -0.05, "Green"),
        ]
        self.generator = PPTPortfolioGenerator(self.projects, "unused.pptx")

    def test_portfolio_copy_uses_current_project_metrics(self):
        takeaways = " ".join(self.generator._portfolio_takeaways())
        themes = " ".join(
            f"{theme['title']} {theme['body']}"
            for theme in self.generator._risk_themes()
        )
        recommendations = " ".join(
            f"{item['title']} {item['text']}"
            for item in self.generator._recommendations()
        )
        generated_copy = f"{takeaways} {themes} {recommendations}"

        self.assertIn("Alpha Delivery", generated_copy)
        self.assertIn("Beta Modernization", generated_copy)
        self.assertNotIn("Titan", generated_copy)
        self.assertNotIn("UniSan", generated_copy)

    def test_generate_creates_seven_slide_deck(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = os.path.join(temp_dir, "monthly.pptx")
            generator = PPTPortfolioGenerator(self.projects, output_path)
            generator.generate()

            presentation = Presentation(output_path)
            self.assertEqual(len(presentation.slides), 7)

            text = " ".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("Alpha Delivery", text)
            self.assertNotIn("UniSan", text)


if __name__ == "__main__":
    unittest.main()
