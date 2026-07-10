import os
import tempfile
import unittest

from pptx import Presentation

from src.weekly_ppt_generator import WeeklyPPTGenerator


class TestWeeklyPPTGenerator(unittest.TestCase):
    def test_generates_readable_four_slide_report(self):
        metrics = {
            "project_name": "Alpha Delivery",
            "project_manager": "Test Manager",
            "today_date": "2026-07-10",
            "rag_status": "Red",
            "calculated_score": 42.0,
            "pct_complete": 0.44,
            "time_elapsed_pct": 0.60,
            "schedule_slippage": 0.16,
            "blockers_count": 2,
            "overdue_milestones_count": 1,
            "project_stage": "Delivery",
            "blockers": [{"task_name": "Client dependency", "status": "Blocked"}],
            "overdue_milestones": [{"task_name": "Integration sign-off"}],
            "reasoning": """### Executive Summary
The project is red because delivery is behind the elapsed timeline.
### Why This Status Was Assigned
Schedule slippage and an active blocker are the main reasons for the status.
### Top Risk Drivers
- Schedule slippage
### Recommended Actions
- Confirm an owner for the client dependency.
- Rebaseline the integration milestone.
### Data Quality & Assumptions
Dates and task statuses were read from the supplied workbook.
""",
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = os.path.join(temp_dir, "Alpha Delivery_weekly_report.pptx")
            WeeklyPPTGenerator(metrics, output_path).generate()
            presentation = Presentation(output_path)
            text = " ".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

        self.assertEqual(len(presentation.slides), 4)
        self.assertIn("Alpha Delivery", text)
        self.assertIn("WHY THIS STATUS", text.upper())
        self.assertIn("RECOMMENDED ACTIONS", text.upper())


if __name__ == "__main__":
    unittest.main()
